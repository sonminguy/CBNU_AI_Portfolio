from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

import common.gridworld_render_Quiz as render_helper
from common.gridworld_Quiz import GridWorld
from mc_control_quiz import MCAgent


ROOT_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = ROOT_DIR / "result" / "quiz_ppt"
EPISODES = 1000
EPSILON_VALUES = [0.05, 0.1, 0.3]
ALPHA_VALUES = [0.01, 0.1, 0.3]
BASE_GAMMA = 0.9


def rolling_mean(values, window=100):
    values = np.asarray(values, dtype=float)
    if len(values) < window:
        return np.arange(len(values)), values
    kernel = np.ones(window) / window
    smoothed = np.convolve(values, kernel, mode="valid")
    xs = np.arange(window - 1, len(values))
    return xs, smoothed


def greedy_policy(agent, env):
    policy = {}
    for state in env.states():
        if state in env.wall_state:
            continue
        qs = [agent.Q[(state, action)] for action in range(agent.action_size)]
        max_q = max(qs)
        probs = {0: 0.0, 1: 0.0, 2: 0.0, 3: 0.0}
        best_actions = [action for action, value in enumerate(qs) if value == max_q]
        share = 1.0 / len(best_actions)
        for action in best_actions:
            probs[action] = share
        policy[state] = probs
    return policy


def policy_text(policy, env):
    arrows = {0: "↑", 1: "↓", 2: "←", 3: "→"}
    lines = []
    for y in range(env.height):
        row = []
        for x in range(env.width):
            state = (y, x)
            if state in env.wall_state:
                row.append("■")
            elif state == env.goal_state:
                row.append("G")
            else:
                best = [action for action, prob in policy[state].items() if prob > 0]
                row.append("".join(arrows[action] for action in best))
        lines.append(" ".join(f"{cell:4s}" for cell in row))
    return "\n".join(lines)


def train_agent(epsilon=0.1, alpha=0.1, episodes=EPISODES):
    env = GridWorld()
    agent = MCAgent()
    agent.gamma = BASE_GAMMA
    agent.epsilon = epsilon
    agent.alpha = alpha
    rewards = []
    steps_per_episode = []
    goal_count = 0

    for _ in range(episodes):
        state = env.reset()
        agent.reset()
        total_reward = 0.0
        step_count = 0

        while True:
            action = agent.get_action(state)
            next_state, reward, done = env.step(action)
            agent.add(state, action, reward)
            total_reward += reward
            step_count += 1
            if done:
                agent.update()
                goal_count += 1
                break
            state = next_state

        rewards.append(total_reward)
        steps_per_episode.append(step_count)

    policy = greedy_policy(agent, env)
    metrics = {
        "epsilon": epsilon,
        "alpha": alpha,
        "avg_reward_last_200": float(np.mean(rewards[-200:])),
        "avg_steps_last_200": float(np.mean(steps_per_episode[-200:])),
        "goal_rate": float(goal_count / episodes),
        "policy_text": policy_text(policy, env),
    }
    return env, agent, rewards, metrics, policy


def save_q_image(env, agent, image_path):
    renderer = render_helper.Renderer(env.reward_map, env.goal_state, env.wall_state)
    renderer.render_q(agent.Q, False)
    plt.gcf().savefig(image_path, bbox_inches="tight", dpi=150)
    plt.close("all")


def save_policy_image(env, policy, image_path):
    renderer = render_helper.Renderer(env.reward_map, env.goal_state, env.wall_state)
    renderer.render_v(None, policy, False)
    plt.gcf().savefig(image_path, bbox_inches="tight", dpi=150)
    plt.close("all")


def save_reward_plot(results, image_path, title):
    plt.figure(figsize=(9, 4.5))
    for label, result in results.items():
        xs, ys = rolling_mean(result["rewards"], window=100)
        plt.plot(xs, ys, linewidth=2, label=label)
    plt.title(title)
    plt.xlabel("Episode")
    plt.ylabel("Rolling mean reward (window=100)")
    plt.legend()
    plt.tight_layout()
    plt.savefig(image_path, dpi=150)
    plt.close()


def build_pair_image(q_path, policy_path, out_path, title):
    q_image = Image.open(q_path).convert("RGB")
    policy_image = Image.open(policy_path).convert("RGB")
    width = q_image.width + policy_image.width + 60
    height = max(q_image.height, policy_image.height) + 90
    canvas = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(canvas)
    draw.text((30, 20), title, fill="black")
    draw.text((30, 50), "Q-function", fill="black")
    draw.text((q_image.width + 60, 50), "Greedy policy", fill="black")
    canvas.paste(q_image, (30, 80))
    canvas.paste(policy_image, (q_image.width + 60, 80))
    canvas.save(out_path)


def run_suite(kind, values):
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    results = {}
    for value in values:
        if kind == "epsilon":
            epsilon = value
            alpha = 0.1
            label = f"epsilon={value}"
            prefix = f"epsilon_{str(value).replace('.', '_')}"
        else:
            epsilon = 0.1
            alpha = value
            label = f"alpha={value}"
            prefix = f"alpha_{str(value).replace('.', '_')}"

        env, agent, rewards, metrics, policy = train_agent(epsilon=epsilon, alpha=alpha)
        q_path = OUTPUT_DIR / f"{prefix}_q.png"
        policy_path = OUTPUT_DIR / f"{prefix}_policy.png"
        pair_path = OUTPUT_DIR / f"{prefix}_pair.png"
        save_q_image(env, agent, q_path)
        save_policy_image(env, policy, policy_path)
        build_pair_image(q_path, policy_path, pair_path, f"{label}, episodes={EPISODES}")

        results[label] = {
            "metrics": metrics,
            "q_path": q_path,
            "policy_path": policy_path,
            "pair_path": pair_path,
            "rewards": rewards,
        }

    plot_path = OUTPUT_DIR / f"{kind}_reward_curve.png"
    save_reward_plot(results, plot_path, f"Monte Carlo Control {kind} comparison")
    return results, plot_path


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(11.2), Inches(1.2))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Monte Carlo Method Quiz"
    p.font.size = Pt(28)
    p.font.bold = True

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(10.5), Inches(1.5))
    sub_frame = sub_box.text_frame
    p = sub_frame.paragraphs[0]
    p.text = "5x5 GridWorld에서 Q 함수와 policy를 구하고, epsilon/alpha 변화에 따른 결과를 비교"
    p.font.size = Pt(18)

    info = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(5.2), Inches(2.0))
    frame = info.text_frame
    for text in [
        "환경: start=(4,0), goal=(0,4)",
        "보상: goal +1, bomb -1 두 곳, wall 두 칸",
        f"실험: episodes={EPISODES}, gamma={BASE_GAMMA}",
    ]:
        p = frame.paragraphs[0] if not frame.paragraphs[0].text else frame.add_paragraph()
        p.text = text
        p.font.size = Pt(18)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(10.2), Inches(5.5))
    frame = body.text_frame
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0


def add_image_slide(prs, title, image_path, caption_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.2), width=Inches(8.0))
    caption = slide.shapes.add_textbox(Inches(8.8), Inches(1.3), Inches(3.0), Inches(4.8))
    frame = caption.text_frame
    for index, line in enumerate(caption_lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def add_three_image_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    positions = [Inches(0.35), Inches(4.1), Inches(7.85)]
    for left, item in zip(positions, items):
        slide.shapes.add_picture(str(item["image"]), left, Inches(1.25), width=Inches(3.35))
        box = slide.shapes.add_textbox(left, Inches(5.85), Inches(3.35), Inches(1.3))
        frame = box.text_frame
        p = frame.paragraphs[0]
        p.text = item["title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p = frame.add_paragraph()
        p.text = item["body"]
        p.font.size = Pt(12)


def add_policy_text_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    positions = [Inches(0.4), Inches(4.1), Inches(7.8)]
    for left, item in zip(positions, items):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.25), Inches(3.2), Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        shape.line.color.rgb = RGBColor(180, 180, 180)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = item["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        for line in item["policy"].splitlines():
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.name = "Consolas"


def build_ppt(epsilon_results, epsilon_plot, alpha_results, alpha_plot):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    best_epsilon_label = max(
        epsilon_results,
        key=lambda label: epsilon_results[label]["metrics"]["avg_reward_last_200"],
    )
    best_alpha_label = max(
        alpha_results,
        key=lambda label: alpha_results[label]["metrics"]["avg_reward_last_200"],
    )
    worst_alpha_steps_label = max(
        alpha_results,
        key=lambda label: alpha_results[label]["metrics"]["avg_steps_last_200"],
    )

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "Source Code Explanation",
        [
            "mc_control_quiz.py는 epsilon-greedy 정책으로 행동을 선택하고, 종료 후 return G를 이용해 상태-행동가치 Q를 갱신한다.",
            "원본 구조를 유지한 채 GridWorld 환경만 Quiz용 5x5로 교체한 형태다.",
            "control에서는 Q가 가장 큰 행동 쪽으로 policy가 점점 이동하고, epsilon이 탐험 강도를 조절한다.",
            "이번 보고서는 control 실험만 포함하고, epsilon/alpha 변화 비교에 집중한다.",
        ],
    )
    add_bullet_slide(
        prs,
        "Experiment Setting",
        [
            f"기본 고정값: gamma={BASE_GAMMA}, episodes={EPISODES}",
            "epsilon 비교: alpha=0.1로 고정하고 epsilon을 0.05, 0.1, 0.3으로 변경",
            "alpha 비교: epsilon=0.1로 고정하고 alpha를 0.01, 0.1, 0.3으로 변경",
            "각 실험에서 reward rolling mean, 최종 Q 함수, greedy policy 이미지를 저장해 비교",
        ],
    )

    add_image_slide(
        prs,
        "Epsilon Comparison: Reward Trend",
        epsilon_plot,
        [
            f"epsilon=0.05: avg reward {epsilon_results['epsilon=0.05']['metrics']['avg_reward_last_200']:.3f}",
            f"epsilon=0.1: avg reward {epsilon_results['epsilon=0.1']['metrics']['avg_reward_last_200']:.3f}",
            f"epsilon=0.3: avg reward {epsilon_results['epsilon=0.3']['metrics']['avg_reward_last_200']:.3f}",
            "이번 실행에서는 epsilon 값에 따라 reward 수준과 수렴 속도가 다르게 나타났다.",
        ],
    )

    add_three_image_slide(
        prs,
        "Epsilon Comparison: Q and Policy Images",
        [
            {
                "image": epsilon_results['epsilon=0.05']['pair_path'],
                "title": "epsilon=0.05",
                "body": f"goal_rate={epsilon_results['epsilon=0.05']['metrics']['goal_rate']:.3f}\navg_steps={epsilon_results['epsilon=0.05']['metrics']['avg_steps_last_200']:.1f}",
            },
            {
                "image": epsilon_results['epsilon=0.1']['pair_path'],
                "title": "epsilon=0.1",
                "body": f"goal_rate={epsilon_results['epsilon=0.1']['metrics']['goal_rate']:.3f}\navg_steps={epsilon_results['epsilon=0.1']['metrics']['avg_steps_last_200']:.1f}",
            },
            {
                "image": epsilon_results['epsilon=0.3']['pair_path'],
                "title": "epsilon=0.3",
                "body": f"goal_rate={epsilon_results['epsilon=0.3']['metrics']['goal_rate']:.3f}\navg_steps={epsilon_results['epsilon=0.3']['metrics']['avg_steps_last_200']:.1f}",
            },
        ],
    )

    add_policy_text_slide(
        prs,
        "Epsilon Comparison: Greedy Policy",
        [
            {"title": "epsilon=0.05", "policy": epsilon_results['epsilon=0.05']['metrics']['policy_text']},
            {"title": "epsilon=0.1", "policy": epsilon_results['epsilon=0.1']['metrics']['policy_text']},
            {"title": "epsilon=0.3", "policy": epsilon_results['epsilon=0.3']['metrics']['policy_text']},
        ],
    )

    add_image_slide(
        prs,
        "Alpha Comparison: Reward Trend",
        alpha_plot,
        [
            f"alpha=0.01: avg reward {alpha_results['alpha=0.01']['metrics']['avg_reward_last_200']:.3f}",
            f"alpha=0.1: avg reward {alpha_results['alpha=0.1']['metrics']['avg_reward_last_200']:.3f}",
            f"alpha=0.3: avg reward {alpha_results['alpha=0.3']['metrics']['avg_reward_last_200']:.3f}",
            "이번 실행에서는 alpha 값에 따라 reward 수준과 평균 이동 횟수 차이가 관찰되었다.",
        ],
    )

    add_three_image_slide(
        prs,
        "Alpha Comparison: Q and Policy Images",
        [
            {
                "image": alpha_results['alpha=0.01']['pair_path'],
                "title": "alpha=0.01",
                "body": f"goal_rate={alpha_results['alpha=0.01']['metrics']['goal_rate']:.3f}\navg_steps={alpha_results['alpha=0.01']['metrics']['avg_steps_last_200']:.1f}",
            },
            {
                "image": alpha_results['alpha=0.1']['pair_path'],
                "title": "alpha=0.1",
                "body": f"goal_rate={alpha_results['alpha=0.1']['metrics']['goal_rate']:.3f}\navg_steps={alpha_results['alpha=0.1']['metrics']['avg_steps_last_200']:.1f}",
            },
            {
                "image": alpha_results['alpha=0.3']['pair_path'],
                "title": "alpha=0.3",
                "body": f"goal_rate={alpha_results['alpha=0.3']['metrics']['goal_rate']:.3f}\navg_steps={alpha_results['alpha=0.3']['metrics']['avg_steps_last_200']:.1f}",
            },
        ],
    )

    add_policy_text_slide(
        prs,
        "Alpha Comparison: Greedy Policy",
        [
            {"title": "alpha=0.01", "policy": alpha_results['alpha=0.01']['metrics']['policy_text']},
            {"title": "alpha=0.1", "policy": alpha_results['alpha=0.1']['metrics']['policy_text']},
            {"title": "alpha=0.3", "policy": alpha_results['alpha=0.3']['metrics']['policy_text']},
        ],
    )

    add_bullet_slide(
        prs,
        "Discussion",
        [
            f"control 실험에서는 이번 실행 기준으로 {best_epsilon_label}이 가장 높은 평균 reward를 보였다.",
            f"alpha 비교에서도 이번 실행 기준으로 {best_alpha_label}이 가장 높은 평균 reward를 보였지만, 이 결과만으로 일반적인 최적값이라고 단정하긴 어렵다.",
            f"{worst_alpha_steps_label}은 평균 이동 횟수가 가장 크게 나타났고, 이는 값 갱신 폭이 커질 때 경로가 불안정해질 수 있음을 시사한다.",
            "epsilon과 alpha는 탐험-활용 균형과 학습 안정성에 직접 영향을 주므로, 동일 환경에서 함께 튜닝하는 접근이 필요하다.",
        ],
    )

    output_path = OUTPUT_DIR / "MonteCarlo_Quiz_Report.pptx"
    try:
        prs.save(output_path)
    except PermissionError:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = OUTPUT_DIR / f"MonteCarlo_Quiz_Report_{timestamp}.pptx"
        prs.save(output_path)
    return output_path


def write_summary(epsilon_results, alpha_results):
    lines = [
        "Monte Carlo Quiz PPT Summary",
        f"episodes={EPISODES}, gamma={BASE_GAMMA}",
        "",
    ]
    lines.extend(
        [
            "[epsilon]",
        ]
    )
    for label, result in epsilon_results.items():
        metrics = result["metrics"]
        lines.extend(
            [
                label,
                f"avg_reward_last_200={metrics['avg_reward_last_200']:.4f}",
                f"avg_steps_last_200={metrics['avg_steps_last_200']:.4f}",
                f"goal_rate={metrics['goal_rate']:.4f}",
                metrics["policy_text"],
                "",
            ]
        )
    lines.append("[alpha]")
    for label, result in alpha_results.items():
        metrics = result["metrics"]
        lines.extend(
            [
                label,
                f"avg_reward_last_200={metrics['avg_reward_last_200']:.4f}",
                f"avg_steps_last_200={metrics['avg_steps_last_200']:.4f}",
                f"goal_rate={metrics['goal_rate']:.4f}",
                metrics["policy_text"],
                "",
            ]
        )
    (OUTPUT_DIR / "summary.txt").write_text("\n".join(lines), encoding="utf-8")


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    epsilon_results, epsilon_plot = run_suite("epsilon", EPSILON_VALUES)
    alpha_results, alpha_plot = run_suite("alpha", ALPHA_VALUES)
    write_summary(epsilon_results, alpha_results)
    ppt_path = build_ppt(epsilon_results, epsilon_plot, alpha_results, alpha_plot)
    print(f"Created PPT: {ppt_path}")


if __name__ == "__main__":
    main()